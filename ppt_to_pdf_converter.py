# -*- coding: utf-8 -*-
"""
Created on Tue Oct 29 14:34:31 2024

@author: Ken
"""

import os
import base64
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from html import escape
from PIL import Image
from io import BytesIO
from xhtml2pdf import pisa
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from xhtml2pdf.default import DEFAULT_FONT
import reportlab.lib.styles

class PPTToPDFConverter:
    
    def __init__(self, pptx_file_path, output_pdf_path, font_path):
        """
        Initialize the PPTToPDFConverter class.
        
        :param pptx_file_path: str, path to the input PPTX file
        :param output_pdf_path: str, path to the output PDF file
        :param font_path: str, path to the font file
        """
        self.pptx_file_path = pptx_file_path
        self.output_pdf_path = output_pdf_path
        self.font_path = font_path
    
    def split_text_to_paragraphs(self, text, max_length=61):
        lines = []
        while len(text) > max_length:
            split_pos = text.rfind(' ', 0, max_length)
            if split_pos == -1:
                split_pos = max_length
            lines.append(text[:split_pos])
            text = text[split_pos:].strip()
        lines.append(text)
        return lines

    def resize_image(self, image_data, max_width, max_height):
        image = Image.open(BytesIO(image_data))
        image.thumbnail((max_width, max_height), Image.LANCZOS)
        buffer = BytesIO()
        image.save(buffer, format="PNG")
        return buffer.getvalue()

    def pptx_to_html(self, output_dir):
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        prs = Presentation(self.pptx_file_path)
        html_output = []

        html_output.append("""
        <head>
            <meta charset="UTF-8">
            <style type="text/css">
                @page {
                    size: A4 landscape;
                    margin: 20mm; 
                }
                body {
                    font-family: 'JapaneseFont'; 
                    font-size: 12pt; 
                    margin: 0; 
                    padding: 0; 
                    white-space: normal;
                    word-wrap: break-word;
                    overflow-x: hidden;
                }
                .slide {
                    width: 257mm; 
                    height: auto; 
                    margin: 0 auto; 
                    padding: 0;
                    overflow: hidden;
                    page-break-inside: avoid;
                    page-break-after: always;
                    display: block;
                    clear: both;
                    position: relative;
                    box-sizing: border-box;
                }
                table {
                    width: 80%; 
                    border-collapse: collapse;
                    margin: 20px auto; 
                    page-break-inside: avoid; 
                }
                th, td {
                    border: 1px solid black;
                    padding: 8px;
                    text-align: left;
                    page-break-inside: avoid;
                }
                p {
                    margin: 0; 
                    padding: 5px 0; 
                }
                img {
                    display: block;
                    margin: 0 auto;
                    max-width: 100%;
                    max-height: 100%;
                    height: auto;
                    width: auto;
                }
            </style>
        </head>
        <body>
        """)

        for slide_number, slide in enumerate(prs.slides):
            html_output.append(f"<div class='slide' id='slide-{slide_number}'>")

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = escape(paragraph.text)
                        split_paragraphs = self.split_text_to_paragraphs(text)
                        for line in split_paragraphs:
                            html_output.append(f"<p>{line}</p>")
                
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_data = shape.image.blob
                    resized_image_data = self.resize_image(image_data, 480, 320)
                    image_base64 = base64.b64encode(resized_image_data).decode('utf-8')
                    html_output.append(f"<img src='data:image/png;base64,{image_base64}' alt='Slide Image' />")

                elif shape.has_table:
                    table = shape.table
                    html_output.append("<table>")
                    for row in table.rows:
                        html_output.append("<tr>")
                        for cell in row.cells:
                            cell_text = escape(cell.text)
                            split_paragraphs = self.split_text_to_paragraphs(cell_text)
                            for line in split_paragraphs:
                                html_output.append(f"<td><p>{line}</p></td>")
                        html_output.append("</tr>")
                    html_output.append("</table>")

            html_output.append("</div>")

        html_output.append("</body>")

        html_content = "\n".join(html_output)
        html_file_path = os.path.join(output_dir, "presentation.html")

        with open(html_file_path, 'w', encoding='utf-8') as html_file:
            html_file.write(html_content)

        return html_file_path

    def convert_html_to_pdf(self, html_path):
        os.makedirs(os.path.dirname(self.output_pdf_path), exist_ok=True)
        pdfmetrics.registerFont(TTFont('JapaneseFont', self.font_path))
        reportlab.lib.styles.ParagraphStyle.defaults['wordWrap'] = 'CJK'
        DEFAULT_FONT['helvetica'] = 'JapaneseFont'

        with open(html_path, 'r', encoding='utf-8') as html_file:
            html_content = html_file.read()

        with open(self.output_pdf_path, 'wb') as pdf_file:
            pisa_status = pisa.CreatePDF(
                html_content,
                dest=pdf_file,
                encoding='utf-8',
            )

        if pisa_status.err:
            print("Error converting HTML to PDF.")
        else:
            print(f"PDF file has been created at {self.output_pdf_path}")

    def convert_ppt_to_pdf(self):
        # Temporary directory to save HTML file
        temp_dir = os.path.dirname(self.output_pdf_path)
        html_path = self.pptx_to_html(temp_dir)
        self.convert_html_to_pdf(html_path)
        print("Conversion from PPTX to PDF completed.")


